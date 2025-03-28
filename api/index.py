from flask import Flask, request, send_file, jsonify
from flask_cors import CORS  # Import CORS
import os
import time
import pandas as pd
from fpdf import FPDF
import nbformat
import tempfile
from docx import Document
from docx.shared import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from PIL import Image, ImageDraw
import io

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# Temporary directory for file uploads and conversions
UPLOAD_FOLDER = tempfile.mkdtemp()

def safe_remove(file_path):
    """Safely remove a file with retries."""
    for _ in range(5):  # Try up to 5 times
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
            break
        except PermissionError:
            time.sleep(0.5)  # Wait and retry

# Convert Word to PDF (Replaced docx2pdf/pypandoc with python-docx + reportlab)
@app.route('/convert/word-to-pdf', methods=['POST'])
def word_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "converted.pdf")
    file.save(input_path)

    try:
        doc = Document(input_path)
        pdf_canvas = canvas.Canvas(output_path, pagesize=letter)
        pdf_canvas.setFont("Helvetica", 12)
        y_position = 750  

        # Process text
        for para in doc.paragraphs:
            pdf_canvas.drawString(100, y_position, para.text)
            y_position -= 20  

            if y_position < 50:
                pdf_canvas.showPage()
                pdf_canvas.setFont("Helvetica", 12)
                y_position = 750

        # Process images
        for rel in doc.part.rels:
            if "image" in doc.part.rels[rel].target_ref:
                img = doc.part.rels[rel].target_part.blob

                img_path = os.path.join(UPLOAD_FOLDER, f"temp_image_{rel}.jpg")
                with open(img_path, "wb") as f:
                    f.write(img)

                img_reader = ImageReader(img_path)
                img_width, img_height = img_reader.getSize()
                aspect_ratio = img_width / img_height

                new_width = 300
                new_height = new_width / aspect_ratio

                if y_position - new_height < 50:  
                    pdf_canvas.showPage()
                    y_position = 750

                pdf_canvas.drawImage(img_reader, 100, y_position - new_height, width=new_width, height=new_height)
                y_position -= new_height + 20

        pdf_canvas.save()
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

    finally:
        safe_remove(input_path)
        safe_remove(output_path)

# Convert Excel to PDF
@app.route('/convert/excel-to-pdf', methods=['POST'])
def excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "converted.pdf")
    file.save(input_path)

    try:
        df = pd.read_excel(input_path)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        for index, row in df.iterrows():
            for col in df.columns:
                pdf.cell(40, 10, str(row[col]), border=1)
            pdf.ln()

        pdf.output(output_path)
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        safe_remove(input_path)
        safe_remove(output_path)

# Convert IPython Notebook to PDF (without pandoc)
@app.route('/convert/ipynb-to-pdf', methods=['POST'])
def ipynb_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "converted.pdf")
    file.save(input_path)

    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            notebook = nbformat.read(f, as_version=4)

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        # Extract notebook cells
        for cell in notebook.cells:
            if cell.cell_type == "markdown":
                pdf.set_font("Arial", style='B', size=14)
                pdf.multi_cell(0, 10, cell.source)  # Markdown text
                pdf.ln(5)
            elif cell.cell_type == "code":
                pdf.set_font("Courier", size=10)
                pdf.multi_cell(0, 8, cell.source)  # Code text
                pdf.ln(5)

        pdf.output(output_path)
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")

    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        safe_remove(input_path)
        safe_remove(output_path)

@app.route('/convert/ppt-to-pdf', methods=['POST'])
def ppt_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No file selected"}), 400

    # Validate file extension
    if not file.filename.lower().endswith(('.ppt', '.pptx')):
        return jsonify({"error": "Invalid file type. Only PPT/PPTX files are allowed."}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "converted.pdf")
    file.save(input_path)

    try:
        # Load PowerPoint file from disk
        with open(input_path, 'rb') as f:
            prs = Presentation(f)

        c = canvas.Canvas(output_path, pagesize=letter)

        for i, slide in enumerate(prs.slides):
            img = Image.new('RGB', (800, 600), 'white')
            draw = ImageDraw.Draw(img)

            draw.text((50, 50), f"Slide {i+1}", fill='black')

            y_offset = 100
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    draw.text((50, y_offset), shape.text, fill='black')
                    y_offset += 30

            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            c.drawImage(ImageReader(img_buffer), 50, 50, width=500, height=400, preserveAspectRatio=True)
            c.showPage()

        c.save()

        return send_file(output_path, as_attachment=True, download_name="converted.pdf")

    except Exception as e:
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500

    finally:
        safe_remove(input_path)
        safe_remove(output_path)

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
