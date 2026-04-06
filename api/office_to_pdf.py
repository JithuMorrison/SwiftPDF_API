from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import os
import time
import tempfile
import subprocess
import platform
import threading
import io

# Basic fallback libs
import pandas as pd
from fpdf import FPDF
from docx import Document
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from pptx import Presentation
from PIL import Image, ImageDraw

app = Flask(__name__)
CORS(app)

UPLOAD_FOLDER = tempfile.mkdtemp()


def safe_remove(file_path):
    """Safely remove a file with retries."""
    for _ in range(5):
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
            break
        except PermissionError:
            time.sleep(0.5)


def convert_with_libreoffice(input_path, output_dir):
    """Convert office file to PDF using LibreOffice (cross-platform)."""
    libreoffice_cmd = None

    if platform.system() == "Windows":
        candidates = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for c in candidates:
            if os.path.exists(c):
                libreoffice_cmd = c
                break
    elif platform.system() == "Darwin":
        libreoffice_cmd = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    else:
        libreoffice_cmd = "libreoffice"

    if libreoffice_cmd is None:
        raise RuntimeError("LibreOffice not found. Please install LibreOffice.")

    result = subprocess.run(
        [libreoffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", output_dir, input_path],
        capture_output=True,
        text=True,
        timeout=60,
    )

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    base_name = os.path.splitext(os.path.basename(input_path))[0]
    output_pdf = os.path.join(output_dir, base_name + ".pdf")

    if not os.path.exists(output_pdf):
        raise RuntimeError("PDF output not found after conversion.")

    return output_pdf


def convert_with_win32com_word(input_path, output_path):
    """Convert Word to PDF using Microsoft Word COM (Windows only)."""
    import win32com.client
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    try:
        doc = word.Documents.Open(os.path.abspath(input_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = wdFormatPDF
        doc.Close()
    finally:
        word.Quit()


def convert_with_win32com_excel(input_path, output_path):
    """Convert Excel to PDF using Microsoft Excel COM (Windows only)."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()  # Required when called from a Flask worker thread
    try:
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        try:
            wb = excel.Workbooks.Open(os.path.abspath(input_path))
            wb.ExportAsFixedFormat(0, os.path.abspath(output_path))  # 0 = xlTypePDF
            wb.Close(False)
        finally:
            excel.Quit()
    finally:
        pythoncom.CoUninitialize()


def convert_with_win32com_ppt(input_path, output_path):
    """Convert PowerPoint to PDF using Microsoft PowerPoint COM (Windows only)."""
    import win32com.client
    import pythoncom
    pythoncom.CoInitialize()  # Required when called from a Flask worker thread
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = 1  # PowerPoint requires Visible=True to open files
        try:
            presentation = ppt_app.Presentations.Open(
                os.path.abspath(input_path), WithWindow=True
            )
            presentation.SaveAs(os.path.abspath(output_path), FileFormat=32)  # 32 = ppSaveAsPDF
            presentation.Close()
        finally:
            ppt_app.Quit()
    finally:
        pythoncom.CoUninitialize()


def basic_word_to_pdf(input_path, output_path):
    """Fallback: convert Word to PDF using python-docx + reportlab."""
    doc = Document(input_path)
    pdf_canvas = canvas.Canvas(output_path, pagesize=letter)
    pdf_canvas.setFont("Helvetica", 12)
    y_position = 750

    for para in doc.paragraphs:
        pdf_canvas.drawString(100, y_position, para.text)
        y_position -= 20
        if y_position < 50:
            pdf_canvas.showPage()
            pdf_canvas.setFont("Helvetica", 12)
            y_position = 750

    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            img_data = doc.part.rels[rel].target_part.blob
            img_path = os.path.join(UPLOAD_FOLDER, f"temp_image_{rel}.jpg")
            with open(img_path, "wb") as f:
                f.write(img_data)
            img_reader = ImageReader(img_path)
            img_width, img_height = img_reader.getSize()
            new_width = 300
            new_height = new_width / (img_width / img_height)
            if y_position - new_height < 50:
                pdf_canvas.showPage()
                y_position = 750
            pdf_canvas.drawImage(img_reader, 100, y_position - new_height, width=new_width, height=new_height)
            y_position -= new_height + 20
            safe_remove(img_path)

    pdf_canvas.save()


def basic_excel_to_pdf(input_path, output_path):
    """Fallback: convert Excel to PDF using pandas + fpdf."""
    df = pd.read_excel(input_path)
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=10)
    for _, row in df.iterrows():
        for col in df.columns:
            pdf.cell(40, 10, str(row[col]), border=1)
        pdf.ln()
    pdf.output(output_path)


def basic_ppt_to_pdf(input_path, output_path):
    """Fallback: convert PPT to PDF using python-pptx + reportlab."""
    with open(input_path, 'rb') as f:
        prs = Presentation(f)
    c = canvas.Canvas(output_path, pagesize=letter)
    for i, slide in enumerate(prs.slides):
        img = Image.new('RGB', (800, 600), 'white')
        draw = ImageDraw.Draw(img)
        draw.text((50, 50), f"Slide {i+1}", fill='black')
        y_offset = 100
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                draw.text((50, y_offset), shape.text, fill='black')
                y_offset += 30
        img_buffer = io.BytesIO()
        img.save(img_buffer, format='PNG')
        img_buffer.seek(0)
        c.drawImage(ImageReader(img_buffer), 50, 50, width=500, height=400, preserveAspectRatio=True)
        c.showPage()
    c.save()


def office_convert(input_path, output_path, app_type):
    """
    3-tier conversion:
      1. win32com (Windows + MS Office) — best fidelity
      2. LibreOffice headless — good fidelity
      3. Basic python libs (python-docx/pptx/pandas) — last resort
    """
    errors = []

    # Tier 1: win32com (Windows only)
    if platform.system() == "Windows":
        try:
            if app_type == "word":
                convert_with_win32com_word(input_path, output_path)
            elif app_type == "excel":
                convert_with_win32com_excel(input_path, output_path)
            elif app_type == "ppt":
                convert_with_win32com_ppt(input_path, output_path)
            return
        except Exception as e:
            errors.append(f"win32com: {e}")

    # Tier 2: LibreOffice
    try:
        out_dir = os.path.dirname(output_path)
        pdf_path = convert_with_libreoffice(input_path, out_dir)
        if pdf_path != output_path:
            os.replace(pdf_path, output_path)
        return
    except Exception as e:
        errors.append(f"LibreOffice: {e}")

    # Tier 3: basic python fallback
    try:
        if app_type == "word":
            basic_word_to_pdf(input_path, output_path)
        elif app_type == "excel":
            basic_excel_to_pdf(input_path, output_path)
        elif app_type == "ppt":
            basic_ppt_to_pdf(input_path, output_path)
        return
    except Exception as e:
        errors.append(f"basic fallback: {e}")

    raise RuntimeError("All conversion methods failed: " + " | ".join(errors))


# ── Word to PDF ──────────────────────────────────────────────────────────────

@app.route('/convert/word-to-pdf', methods=['POST'])
def word_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({"error": "No file selected"}), 400
    if not file.filename.lower().endswith(('.doc', '.docx')):
        return jsonify({"error": "Invalid file type. Only DOC/DOCX allowed."}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "word_converted.pdf")
    file.save(input_path)

    try:
        office_convert(input_path, output_path, "word")
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        safe_remove(input_path)
        safe_remove(output_path)


# ── Excel to PDF ─────────────────────────────────────────────────────────────

@app.route('/convert/excel-to-pdf', methods=['POST'])
def excel_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({"error": "No file selected"}), 400
    if not file.filename.lower().endswith(('.xls', '.xlsx')):
        return jsonify({"error": "Invalid file type. Only XLS/XLSX allowed."}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "excel_converted.pdf")
    file.save(input_path)

    try:
        office_convert(input_path, output_path, "excel")
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500
    finally:
        safe_remove(input_path)
        safe_remove(output_path)


# ── PowerPoint to PDF ────────────────────────────────────────────────────────

@app.route('/convert/ppt-to-pdf', methods=['POST'])
def ppt_to_pdf():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({"error": "No file selected"}), 400
    if not file.filename.lower().endswith(('.ppt', '.pptx')):
        return jsonify({"error": "Invalid file type. Only PPT/PPTX allowed."}), 400

    input_path = os.path.join(UPLOAD_FOLDER, file.filename)
    output_path = os.path.join(UPLOAD_FOLDER, "ppt_converted.pdf")
    file.save(input_path)

    try:
        office_convert(input_path, output_path, "ppt")
        return send_file(output_path, as_attachment=True, download_name="converted.pdf")
    except Exception as e:
        return jsonify({"error": f"Conversion failed: {str(e)}"}), 500
    finally:
        safe_remove(input_path)
        safe_remove(output_path)


@app.route('/shutdown', methods=['GET', 'POST'])
def shutdown():
    """Gracefully shut down the Flask server."""
    def _stop():
        time.sleep(0.5)
        os._exit(0)
    threading.Thread(target=_stop, daemon=True).start()
    return jsonify({"message": "Server shutting down..."}), 200


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5001, debug=False, use_reloader=False)
